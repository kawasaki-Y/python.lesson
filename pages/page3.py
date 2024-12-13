import streamlit as st
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# ページ設定
st.set_page_config(page_title="事業計画とシミュレーション", page_icon="", layout="wide")

# CSS
st.markdown(
    """
    <style>
    body {
        background-color: #f0f8ff;
    }
    .stDataFrame, .stMetric {
        font-size: 16px;
    }
    .block-container {
        padding-top: 1rem;
        padding-bottom: 1rem;
        padding-left: 2rem;
        padding-right: 2rem;
    }
    .sidebar .sidebar-content {
        background: #f8f9fa;
    }
    .stPlotlyChart, .stPyplot {
        margin: auto;
        border: 1px solid #ddd;
        border-radius: 10px;
        padding: 10px;
        background-color: #fff;
        box-shadow: 0px 4px 6px rgba(0, 0, 0, 0.1);
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ページタイトル
st.title(" 事業計画とシミュレーション")

# サイドバー
st.sidebar.header("入力情報")
num_businesses = st.sidebar.slider("事業の数", min_value=1, max_value=5, value=2)
simulation_years = st.sidebar.slider("シミュレーション年数", min_value=1, max_value=10, value=5)
currency_unit = st.sidebar.selectbox("通貨単位", ["万円", "千円", "ドル"])
unit_factor = {"万円": 1, "千円": 0.1, "ドル": 0.007}[currency_unit]

# 追加の財務情報入力
st.sidebar.header("財務指標")
total_assets = st.sidebar.number_input("総資産 (万円)", min_value=1, step=10, value=1000) * unit_factor

# サイドバーで事業ごとの年間売上・コスト入力
st.sidebar.header("事業ごとの年間売上・コストと成長率")
business_data = []
for i in range(num_businesses):
    st.sidebar.subheader(f"事業 {i+1}")
    revenue = st.sidebar.number_input(f"事業 {i+1} の年間売上予測 ({currency_unit})", min_value=0, step=10, key=f"revenue_{i}") * unit_factor
    cost = st.sidebar.number_input(f"事業 {i+1} の年間コスト予測 ({currency_unit})", min_value=0, step=10, key=f"cost_{i}") * unit_factor
    growth_rate = st.sidebar.slider(f"事業 {i+1} の年間成長率 (%)", min_value=0, max_value=50, value=10, key=f"growth_rate_{i}")
    business_data.append({"事業": f"事業 {i+1}", "売上": revenue, "コスト": cost, "成長率": growth_rate})

# メインエリアでシミュレーション結果を表示
st.header(" 事業計画シミュレーション")

# シミュレーションデータを計算
simulation_data = []
for year in range(1, simulation_years + 1):
    year_data = {"年数": f"{year}年目"}
    total_profit = 0
    for business in business_data:
        annual_revenue = business["売上"] * (1 + business["成長率"] / 100) ** (year - 1)
        annual_cost = business["コスト"] * (1 + business["成長率"] / 100) ** (year - 1)
        annual_profit = annual_revenue - annual_cost
        year_data[business["事業"]] = round(annual_profit)
        total_profit += annual_profit
    year_data["営業利益合計"] = round(total_profit)
    year_data["ROA"] = round(total_profit / total_assets * 100, 2)  # 総資産利益率
    simulation_data.append(year_data)

# データフレーム作成
simulation_df = pd.DataFrame(simulation_data)

# シミュレーション結果の表を表示
st.subheader(" シミュレーション結果")
st.dataframe(simulation_df)

# 年数を整数型に変換
simulation_df["年数"] = simulation_df["年数"].str.replace("年目", "").astype(int)

# 営業利益推移のグラフ
st.subheader(" 営業利益推移")
fig_profit, ax_profit = plt.subplots(figsize=(5, 3))
ax_profit.plot(simulation_df["年数"], simulation_df["営業利益合計"], marker="o", label="営業利益合計", color="#1f77b4")
ax_profit.set_title("営業利益推移", fontsize=12)
ax_profit.set_xlabel("年数", fontsize=10)  # X軸ラベル
ax_profit.set_ylabel(f"営業利益額 ({currency_unit})", fontsize=10)  # Y軸ラベル
ax_profit.tick_params(axis='both', labelsize=8)
ax_profit.legend()
ax_profit.grid(True, linestyle="--", alpha=0.7)
fig_profit.tight_layout()  # レイアウトを調整
st.pyplot(fig_profit)

# ROAの折れ線グラフ
st.subheader(" ROAの推移")
fig_roa, ax_roa = plt.subplots(figsize=(5, 3))
ax_roa.plot(simulation_df["年数"], simulation_df["ROA"], marker="o", color="green", label="ROA (%)")
ax_roa.set_xlabel("年数", fontsize=10)  # X軸ラベル
ax_roa.set_ylabel("ROA (%)", fontsize=10)  # Y軸ラベルを%表記
ax_roa.tick_params(axis='both', labelsize=8)
ax_roa.grid(True, linestyle="--", alpha=0.7)
fig_roa.tight_layout()  # レイアウトを調整
st.pyplot(fig_roa)

# Excelダウンロード
st.subheader("ファイルをダウンロード")
excel_file = BytesIO()
with pd.ExcelWriter(excel_file, engine="xlsxwriter") as writer:
    simulation_df.to_excel(writer, index=False, sheet_name="シミュレーション結果")
excel_file.seek(0)
st.download_button(
    label="📤Excelをダウンロード",
    data=excel_file,
    file_name="simulation_results.xlsx",
    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
)

# PowerPointダウンロード
def create_pptx_with_graph(dataframe, fig_profit, fig_roa):
    prs = Presentation()

    # 表スライド
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "事業計画シミュレーション"

    # データを表に挿入
    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table

    # ヘッダー
    for col_idx, col_name in enumerate(dataframe.columns):
        table.cell(0, col_idx).text = col_name

    # データ
    for row_idx, row in dataframe.iterrows():
        for col_idx, cell in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(cell)

    # グラフスライド (利益)
    graph_slide_profit = prs.slides.add_slide(prs.slide_layouts[5])
    graph_title_profit = graph_slide_profit.shapes.title
    graph_title_profit.text = "利益合計の推移"
    image_stream_profit = BytesIO()
    fig_profit.savefig(image_stream_profit, format="png")
    image_stream_profit.seek(0)
    graph_slide_profit.shapes.add_picture(image_stream_profit, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5))

    # グラフスライド (ROA)
    graph_slide_roa = prs.slides.add_slide(prs.slide_layouts[5])
    graph_title_roa = graph_slide_roa.shapes.title
    graph_title_roa.text = "ROAの推移"
    image_stream_roa = BytesIO()
    fig_roa.savefig(image_stream_roa, format="png")
    image_stream_roa.seek(0)
    graph_slide_roa.shapes.add_picture(image_stream_roa, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5))

    return prs

if st.button("📤 PowerPointをダウンロード"):
    pptx_file = BytesIO()
    pptx = create_pptx_with_graph(simulation_df, fig_profit, fig_roa)
    pptx.save(pptx_file)
    pptx_file.seek(0)
    st.download_button(
        label="PowerPointをダウンロード",
        data=pptx_file,
        file_name="business_simulation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

