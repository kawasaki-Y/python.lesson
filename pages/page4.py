import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
from io import BytesIO

# ページ設定
st.set_page_config(page_title="Page3: 財務計画", page_icon="📉", layout="wide")

# ページタイトル
st.title("📉 財務計画とシミュレーション")

# サイドバーで事業数を調整
st.sidebar.header("事業の設定")
num_businesses = st.sidebar.slider("事業の数を選択してください", min_value=1, max_value=5, value=2)

# 事業ごとのデータ入力フォーム
st.header("🔢 事業ごとの年間売上・コスト入力")
business_data = []
for i in range(num_businesses):
    st.subheader(f"事業 {i+1}")
    revenue = st.number_input(f"事業 {i+1} の年間売上予測 (万円)", min_value=0, step=10, key=f"revenue_{i}")
    cost = st.number_input(f"事業 {i+1} の年間コスト予測 (万円)", min_value=0, step=10, key=f"cost_{i}")
    profit = revenue - cost
    st.metric(f"事業 {i+1} の予測利益", f"{profit} 万円")
    business_data.append({"事業": f"事業 {i+1}", "売上": revenue, "コスト": cost, "利益": profit})

# 5カ年計画のシミュレーション
st.header("📈 5カ年の事業計画シミュレーション")

# 各事業の年間成長率を設定
growth_rate = st.slider("年間成長率 (%)", min_value=0, max_value=50, value=10)

# データフレームの作成
df_list = []
for business in business_data:
    revenue = business["売上"]
    cost = business["コスト"]
    business_plan = {"年度": [], "事業": [], "売上": [], "コスト": [], "利益": []}
    for year in range(1, 6):
        annual_revenue = revenue * (1 + growth_rate / 100) ** (year - 1)
        annual_cost = cost * (1 + growth_rate / 100) ** (year - 1)
        annual_profit = annual_revenue - annual_cost
        business_plan["年度"].append(f"{year}年目")
        business_plan["事業"].append(business["事業"])
        business_plan["売上"].append(round(annual_revenue))
        business_plan["コスト"].append(round(annual_cost))
        business_plan["利益"].append(round(annual_profit))
    df_list.append(pd.DataFrame(business_plan))

# データフレームを統合して表示
final_df = pd.concat(df_list, ignore_index=True)
st.dataframe(final_df)

# 利益の合計を年度ごとに集計
total_profit_by_year = (
    final_df.groupby("年度")["利益"]
    .sum()
    .reset_index()
    .rename(columns={"利益": "総利益"})
)

# 利益合計の折れ線グラフ
st.subheader("📊 総利益の推移")
st.line_chart(total_profit_by_year.set_index("年度"))

# PowerPointダウンロード
def create_pptx(dataframe):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "5カ年事業計画シミュレーション"

    # データを表に挿入
    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2)).table

    # ヘッダー
    for col_idx, col_name in enumerate(dataframe.columns):
        table.cell(0, col_idx).text = col_name

    # データ
    for row_idx, row in dataframe.iterrows():
        for col_idx, cell in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(cell)

    return prs

if st.button("📤 PowerPointをダウンロード"):
    pptx_file = BytesIO()
    pptx = create_pptx(final_df)
    pptx.save(pptx_file)
    pptx_file.seek(0)
    st.download_button(
        label="PowerPointをダウンロード",
        data=pptx_file,
        file_name="business_plan.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# PDFダウンロード
def create_pdf(dataframe):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="5カ年事業計画シミュレーション", ln=True, align="C")

    # データを表形式で追加
    pdf.set_font("Arial", size=10)
    col_width = pdf.w / (len(dataframe.columns) + 1)
    row_height = pdf.font_size * 1.5

    # ヘッダー
    for col_name in dataframe.columns:
        pdf.cell(col_width, row_height, txt=col_name, border=1)
    pdf.ln(row_height)

    # データ
    for _, row in dataframe.iterrows():
        for cell in row:
            pdf.cell(col_width, row_height, txt=str(cell), border=1)
        pdf.ln(row_height)

    return pdf

if st.button("📥 PDFをダウンロード"):
    pdf_file = BytesIO()
    pdf = create_pdf(final_df)
    pdf.output(pdf_file)
    pdf_file.seek(0)
    st.download_button(
        label="PDFをダウンロード",
        data=pdf_file,
        file_name="business_plan.pdf",
        mime="application/pdf",
    )
